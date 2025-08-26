#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import PyPDF2
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import os

def extract_pdf_text(pdf_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            num_pages = len(pdf_reader.pages)
            
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"
                
        return text, num_pages
    except Exception as e:
        print(f"Error reading {pdf_path}: {str(e)}")
        return None, 0

def parse_korean_pdf():
    """Parse the Korean SDV document"""
    pdf_path = "TalkFile_SDV 개념 및 중독일 표준화 동향_최동근작성.pdf.pdf"
    text, pages = extract_pdf_text(pdf_path)
    
    if text:
        sections = []
        
        # Try to find major sections based on common patterns
        lines = text.split('\n')
        current_section = {"title": "SDV 개념 및 표준화 동향", "content": []}
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Look for numbered sections or major headings
            if re.match(r'^\d+\.', line) or re.match(r'^[IVX]+\.', line):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"title": line, "content": []}
            else:
                current_section["content"].append(line)
        
        if current_section["content"]:
            sections.append(current_section)
            
        return {
            "title": "SDV 개념 및 중독일 표준화 동향",
            "author": "최동근",
            "pages": pages,
            "sections": sections
        }
    return None

def parse_chinese_pdfs():
    """Try to parse Chinese PDFs"""
    chinese_docs = []
    
    pdf_files = [
        ("SDV Intelligent Connected Vehicle Service Interface Specification Part 1 Atomic Service API Interface Version 4 Beta 1(중국어).pdf", "Part 1: Atomic Service API"),
        ("SDV Intelligent Connected Vehicle Service Interface Specification Part 2 Device Abstraction API Interface Version 4 Beta 1(중국어).pdf", "Part 2: Device Abstraction API")
    ]
    
    for pdf_file, title in pdf_files:
        text, pages = extract_pdf_text(pdf_file)
        if text:
            # Extract what we can from the structure
            chinese_docs.append({
                "title": title,
                "pages": pages,
                "content_preview": text[:500] if text else "Unable to extract content",
                "full_title": pdf_file.replace("(중국어).pdf", "")
            })
    
    return chinese_docs

def create_ppt_presentation(korean_data, chinese_data):
    """Create PowerPoint presentation"""
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SDV (Software-Defined Vehicle)"
    subtitle.text = "개념 및 표준화 동향 분석\n중독일 표준 문서 리뷰"
    
    # Overview Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "문서 개요"
    
    overview_text = "포함된 문서:\n\n"
    overview_text += f"1. SDV 개념 및 중독일 표준화 동향 (작성: 최동근)\n"
    overview_text += f"   - 총 {korean_data['pages']}페이지\n\n"
    overview_text += "2. SDV Intelligent Connected Vehicle Service Interface Specification\n"
    for doc in chinese_data:
        overview_text += f"   - {doc['title']}: {doc['pages']}페이지\n"
    
    content.text = overview_text
    
    # Korean Document Section
    if korean_data and korean_data.get('sections'):
        # Main title slide for Korean document
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        title = slide.shapes.title
        title.text = korean_data['title']
        
        # Add slides for major sections
        for i, section in enumerate(korean_data['sections'][:10]):  # Limit to first 10 sections
            if len(section['content']) > 0:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = section['title'][:100]  # Limit title length
                
                # Combine content, limiting to reasonable amount
                content_text = '\n'.join(section['content'][:10])
                if len(content_text) > 500:
                    content_text = content_text[:500] + "..."
                content.text = content_text
    
    # Chinese Documents Section
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "SDV Service Interface Specifications"
    
    for doc in chinese_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = doc['title']
        content_text = f"Document: {doc['full_title']}\n\n"
        content_text += f"Total Pages: {doc['pages']}\n\n"
        content_text += "Content Preview:\n"
        content_text += doc['content_preview']
        content.text = content_text
    
    # Summary Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "요약 및 결론"
    
    summary_text = "주요 내용:\n\n"
    summary_text += "• SDV는 소프트웨어 중심의 차량 아키텍처\n"
    summary_text += "• 중국, 독일, 일본의 표준화 동향 분석\n"
    summary_text += "• Service Interface Specification 정의\n"
    summary_text += "  - Part 1: Atomic Service API\n"
    summary_text += "  - Part 2: Device Abstraction API\n"
    summary_text += "• Version 4 Beta 1 사양 문서화"
    
    content.text = summary_text
    
    # Save presentation
    prs.save('SDV_Presentation.pptx')
    print("Presentation saved as SDV_Presentation.pptx")

def main():
    print("Starting PDF extraction and PPT creation...")
    
    # Change to the sdv-c directory
    os.chdir('/home/kim/sdv-c')
    
    # Parse Korean PDF
    print("Parsing Korean PDF...")
    korean_data = parse_korean_pdf()
    
    # Parse Chinese PDFs
    print("Attempting to parse Chinese PDFs...")
    chinese_data = parse_chinese_pdfs()
    
    # Create PowerPoint presentation
    print("Creating PowerPoint presentation...")
    create_ppt_presentation(korean_data, chinese_data)
    
    print("Done!")

if __name__ == "__main__":
    main()