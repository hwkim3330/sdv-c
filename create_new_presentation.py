#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

def create_sdv_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Color scheme
    primary_color = RGBColor(0, 84, 159)  # Dark blue
    secondary_color = RGBColor(0, 176, 240)  # Light blue
    accent_color = RGBColor(255, 192, 0)  # Yellow
    text_color = RGBColor(51, 51, 51)  # Dark gray
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Software-Defined Vehicle (SDV)"
    subtitle.text = "글로벌 표준화 동향 및 한국 대응 전략\n2025년 1월"
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = "SDV 시장 전망"
    p = tf.add_paragraph()
    p.text = "• 2030년까지 글로벌 SDV 시장 1,500억 달러 예상"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 소프트웨어가 차량 가치의 60% 이상 차지"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n핵심 기술 트렌드"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 서비스 지향 아키텍처(SOA) 채택 확대"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• OTA 업데이트 표준화"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 클라우드-엣지 컴퓨팅 통합"
    p.level = 1
    
    # Slide 3: SDV 개념 정의
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV(Software-Defined Vehicle) 개념"
    tf = body_shape.text_frame
    tf.text = "정의"
    p = tf.add_paragraph()
    p.text = "• 소프트웨어 중심으로 기능이 정의되고 제어되는 차량"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 하드웨어와 소프트웨어의 분리(Decoupling)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n핵심 특징"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 소프트웨어 업데이트를 통한 지속적 기능 향상"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 표준화된 인터페이스 기반 모듈식 아키텍처"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 실시간 데이터 처리 및 클라우드 연동"
    p.level = 1
    
    # Slide 4: 글로벌 표준화 현황
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "글로벌 SDV 표준화 현황"
    tf = body_shape.text_frame
    tf.text = "중국 - C-ICVS (China Intelligent Connected Vehicle Service)"
    p = tf.add_paragraph()
    p.text = "• 원자 서비스 API 표준 (Version 4 Beta)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 디바이스 추상화 API 표준"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n독일 - AUTOSAR Adaptive Platform"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 동적 소프트웨어 구성 지원"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• POSIX 운영체제 기반"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n일본 - 차량 소프트웨어 플랫폼"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 안전성 중심 설계"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• ISO 26262 준수"
    p.level = 1
    
    # Slide 5: 중국 SDV 표준 상세
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "중국 SDV 표준 심층 분석"
    tf = body_shape.text_frame
    tf.text = "표준화 전략"
    p = tf.add_paragraph()
    p.text = "• 정부 주도 통합 표준 개발"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 산업체 참여 의무화"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nAPI 구조"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 3계층 아키텍처: Application - Service - Hardware"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• RESTful API 기반 통신"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 실시간 이벤트 처리 지원"
    p.level = 1
    
    # Slide 6: 기술 아키텍처
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 기술 아키텍처"
    tf = body_shape.text_frame
    tf.text = "하드웨어 계층"
    p = tf.add_paragraph()
    p.text = "• 고성능 컴퓨팅 플랫폼 (HPC)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 도메인 컨트롤러"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n미들웨어 계층"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 서비스 지향 미들웨어 (SOME/IP)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 데이터 분산 서비스 (DDS)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n애플리케이션 계층"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 차량 서비스 애플리케이션"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 사용자 경험 애플리케이션"
    p.level = 1
    
    # Slide 7: 핵심 기술 요소
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 핵심 기술 요소"
    tf = body_shape.text_frame
    tf.text = "OTA (Over-The-Air) 업데이트"
    p = tf.add_paragraph()
    p.text = "• 차분 업데이트 기술"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 보안 서명 및 검증"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 롤백 메커니즘"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n서비스 오케스트레이션"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 마이크로서비스 관리"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 동적 리소스 할당"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n데이터 관리"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 실시간 데이터 스트리밍"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 엣지-클라우드 동기화"
    p.level = 1
    
    # Slide 8: 보안 아키텍처
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 보안 아키텍처"
    tf = body_shape.text_frame
    tf.text = "보안 위협"
    p = tf.add_paragraph()
    p.text = "• 원격 해킹 및 제어권 탈취"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 데이터 유출 및 프라이버시 침해"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 서비스 거부 공격(DoS)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n보안 대책"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• Hardware Security Module (HSM)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Secure Boot 및 신뢰 체인"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 침입 탐지 시스템 (IDS)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 암호화 통신 (TLS/DTLS)"
    p.level = 1
    
    # Slide 9: 한국 산업 현황
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "한국 SDV 산업 현황"
    tf = body_shape.text_frame
    tf.text = "강점"
    p = tf.add_paragraph()
    p.text = "• 우수한 IT 인프라 및 5G 네트워크"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 반도체 및 디스플레이 기술력"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 완성차-부품사 수직계열화"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n약점"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 소프트웨어 플랫폼 경쟁력 부족"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 글로벌 표준 주도권 미흡"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 전문 인력 부족"
    p.level = 1
    
    # Slide 10: 한국 대응 전략
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "한국 SDV 대응 전략"
    tf = body_shape.text_frame
    tf.text = "단기 전략 (1-2년)"
    p = tf.add_paragraph()
    p.text = "• SDV 전문 인력 양성 프로그램 확대"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 국제 표준화 기구 참여 강화"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 산학연 협력체계 구축"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n중기 전략 (3-5년)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• K-SDV 플랫폼 개발"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 테스트베드 구축 및 실증"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 글로벌 파트너십 확대"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n장기 전략 (5년+)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 차세대 SDV 표준 주도"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• SDV 생태계 글로벌 확장"
    p.level = 1
    
    # Slide 11: 정책 제언
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "정책 제언"
    tf = body_shape.text_frame
    tf.text = "정부 차원"
    p = tf.add_paragraph()
    p.text = "• SDV 특별법 제정"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• R&D 투자 확대 (연간 1조원 규모)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 규제 샌드박스 활성화"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n산업 차원"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 오픈소스 SDV 플랫폼 공동 개발"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 크로스 도메인 협업 강화"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 스타트업 육성 및 M&A 활성화"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n학계 차원"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• SDV 전문 교육과정 신설"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 산학협력 프로젝트 확대"
    p.level = 1
    
    # Slide 12: 시장 기회
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 시장 기회"
    tf = body_shape.text_frame
    tf.text = "새로운 비즈니스 모델"
    p = tf.add_paragraph()
    p.text = "• Features-as-a-Service (FaaS)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 구독형 차량 기능"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 데이터 기반 서비스"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n수익 창출 영역"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• OTA 업데이트 플랫폼"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 차량 데이터 분석"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 개인화 서비스"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 써드파티 앱 마켓플레이스"
    p.level = 1
    
    # Slide 13: 도전 과제
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 구현 도전 과제"
    tf = body_shape.text_frame
    tf.text = "기술적 도전"
    p = tf.add_paragraph()
    p.text = "• 실시간 처리 요구사항 충족"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 기능 안전성 보장"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 레거시 시스템 통합"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n비즈니스 도전"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 높은 초기 투자 비용"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 수익 모델 불확실성"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 파트너십 구축 복잡성"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n규제 도전"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 국가별 상이한 규제"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 데이터 주권 이슈"
    p.level = 1
    
    # Slide 14: 성공 요인
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 성공 요인"
    tf = body_shape.text_frame
    tf.text = "핵심 성공 요인 (Critical Success Factors)"
    p = tf.add_paragraph()
    p.text = "• 표준화된 플랫폼 확보"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 강력한 생태계 구축"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 지속적인 혁신 역량"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 사용자 신뢰 확보"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n경쟁 우위 확보 방안"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 차별화된 사용자 경험"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 빠른 개발 및 배포 사이클"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 데이터 기반 의사결정"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 글로벌 협력 네트워크"
    p.level = 1
    
    # Slide 15: 로드맵
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "SDV 구현 로드맵"
    tf = body_shape.text_frame
    tf.text = "2025-2026: 기반 구축"
    p = tf.add_paragraph()
    p.text = "• 아키텍처 설계 및 표준화"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 파일럿 프로젝트 실행"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n2027-2028: 확산"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 상용화 모델 출시"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 서비스 플랫폼 운영"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n2029-2030: 성숙"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 완전 자율 SDV 실현"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 글로벌 시장 리더십"
    p.level = 1
    
    # Slide 16: 투자 계획
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "투자 계획"
    tf = body_shape.text_frame
    tf.text = "투자 규모 (5년간)"
    p = tf.add_paragraph()
    p.text = "• 플랫폼 개발: 3,000억원"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 인프라 구축: 2,000억원"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 인력 양성: 1,000억원"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 표준화 활동: 500억원"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n기대 효과"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 신규 일자리 10,000개 창출"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 수출 증대 연간 10조원"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 글로벌 시장 점유율 15%"
    p.level = 1
    
    # Slide 17: 글로벌 협력
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "글로벌 협력 전략"
    tf = body_shape.text_frame
    tf.text = "국제 표준화 기구 참여"
    p = tf.add_paragraph()
    p.text = "• ISO/SAE 21434 (사이버보안)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• ISO 26262 (기능안전)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AUTOSAR Consortium"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n전략적 파트너십"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• 글로벌 OEM과 공동 개발"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 빅테크 기업과 기술 협력"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 오픈소스 커뮤니티 참여"
    p.level = 1
    
    # Slide 18: 성과 지표
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "성과 지표 (KPI)"
    tf = body_shape.text_frame
    tf.text = "기술 지표"
    p = tf.add_paragraph()
    p.text = "• OTA 업데이트 성공률 > 99.9%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 서비스 가용성 > 99.95%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 보안 사고 Zero"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n비즈니스 지표"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• SDV 탑재 차량 비율 > 80%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 서비스 구독률 > 60%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 고객 만족도 > 4.5/5.0"
    p.level = 1
    
    # Slide 19: 결론
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "결론 및 제언"
    tf = body_shape.text_frame
    tf.text = "핵심 메시지"
    p = tf.add_paragraph()
    p.text = "• SDV는 자동차 산업의 미래이자 필수 전환점"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 한국은 강점을 활용한 차별화 전략 필요"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 정부-산업-학계 협력이 성공의 열쇠"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n즉시 실행 과제"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• SDV 추진 TF 구성"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 표준화 로드맵 수립"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 파일럿 프로젝트 착수"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 국제 협력 채널 구축"
    p.level = 1
    
    # Slide 20: Q&A
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Q&A"
    subtitle.text = "감사합니다\n\n문의사항: sdv-korea@example.com"
    
    # Save the presentation
    prs.save('/home/kim/sdv-c/SDV_Comprehensive_Presentation_2025.pptx')
    print("Presentation created successfully: SDV_Comprehensive_Presentation_2025.pptx")

if __name__ == "__main__":
    create_sdv_presentation()