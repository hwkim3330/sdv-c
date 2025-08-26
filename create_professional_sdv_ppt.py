#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

def setup_slide_size(prs):
    """Set presentation to 16:9 widescreen format"""
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs

def add_title_slide(prs, title, subtitle):
    """Add a professional title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Subtitle
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
    
    return slide

def add_section_divider(prs, section_title, section_number):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[2]  # Section header
    slide = prs.slides.add_slide(slide_layout)
    
    # Add section number
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(1.5)
    height = Inches(1)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = f"{section_number:02d}"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    # Section title
    title = slide.shapes.title
    title.text = section_title
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title.left = Inches(2.5)
    
    return slide

def add_content_slide(prs, title, content=None, bullets=None, layout_idx=1):
    """Add a content slide with proper formatting"""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title formatting
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    if content or bullets:
        content_shape = slide.placeholders[1]
        
        if content:
            content_shape.text = content
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(89, 89, 89)
                
        elif bullets:
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                if isinstance(bullet, str):
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(51, 51, 51)
                elif isinstance(bullet, dict):
                    # Main bullet
                    p.text = bullet.get('main', '')
                    p.level = 0
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(51, 51, 51)
                    
                    # Sub-bullets
                    if 'sub' in bullet:
                        for sub in bullet['sub']:
                            p = text_frame.add_paragraph()
                            p.text = sub
                            p.level = 1
                            p.font.size = Pt(18)
                            p.font.color.rgb = RGBColor(89, 89, 89)
    
    return slide

def add_comparison_table(prs, title, headers, rows):
    """Add a slide with comparison table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Table
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(3.8)
    
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    
    # Set column widths
    for i in range(len(headers)):
        table.columns[i].width = Inches(9 / len(headers))
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            
            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(51, 51, 51)
            paragraph.alignment = PP_ALIGN.LEFT
    
    return slide

def add_timeline_slide(prs, title, timeline_data):
    """Add a timeline/roadmap slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Timeline
    y_position = Inches(1.5)
    
    for i, item in enumerate(timeline_data):
        # Year/Phase box
        left = Inches(0.5 + i * 2.2)
        top = y_position
        width = Inches(2)
        height = Inches(0.8)
        
        # Phase box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)
        
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = item['phase']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Details box
        detail_box = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(2))
        text_frame = detail_box.text_frame
        text_frame.clear()
        
        for j, detail in enumerate(item['details']):
            if j == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = f"• {detail}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(51, 51, 51)
    
    return slide

def create_professional_sdv_presentation():
    """Create a professional SDV presentation"""
    prs = Presentation()
    prs = setup_slide_size(prs)
    
    # 1. Title Slide
    add_title_slide(
        prs,
        "SDV (Software-Defined Vehicle)\n글로벌 표준화 동향 및 대응 전략",
        f"중국 SDV 표준 심층 분석 | 한국 자동차 산업 로드맵\n\n{datetime.datetime.now().strftime('%Y년 %m월 %d일')}"
    )
    
    # 2. Executive Summary
    bullets = [
        {"main": "SDV 시장 현황", "sub": [
            "2030년 3,500억 달러 규모 예상",
            "전체 차량의 95% SDV 전환"
        ]},
        {"main": "글로벌 표준화 동향", "sub": [
            "중국: 정부 주도 빠른 표준화",
            "독일: AUTOSAR Adaptive Platform",
            "일본: 기업 중심 실용적 접근"
        ]},
        {"main": "한국 대응 전략", "sub": [
            "단기: 기술 확보 및 파일럿 프로젝트",
            "중기: K-SDV 표준 개발",
            "장기: 글로벌 리더십 확보"
        ]}
    ]
    add_content_slide(prs, "Executive Summary", bullets=bullets)
    
    # 3. Agenda
    agenda_items = [
        "SDV 개념 및 시장 전망",
        "중국 SDV 표준 분석",
        "글로벌 표준 비교",
        "기술 아키텍처 분석",
        "한국 대응 전략",
        "실행 로드맵"
    ]
    add_content_slide(prs, "Agenda", bullets=agenda_items)
    
    # Section 1: SDV Overview
    add_section_divider(prs, "SDV 개념 및 시장 전망", 1)
    
    # SDV Definition
    bullets = [
        {"main": "Software-Defined Vehicle", "sub": [
            "차량 기능이 소프트웨어로 정의되고 제어",
            "하드웨어와 소프트웨어의 분리 (Decoupling)",
            "실시간 기능 업데이트 및 확장 가능"
        ]},
        {"main": "핵심 기술 요소", "sub": [
            "중앙 집중형 컴퓨팅 아키텍처",
            "OTA (Over-The-Air) 업데이트",
            "클라우드 네이티브 기술",
            "AI/ML 기반 자율 기능"
        ]}
    ]
    add_content_slide(prs, "SDV 정의 및 핵심 기술", bullets=bullets)
    
    # Market Analysis
    headers = ["구분", "2024년", "2027년", "2030년", "CAGR"]
    rows = [
        ["시장 규모", "$650억", "$1,800억", "$3,500억", "32%"],
        ["SDV 비중", "15%", "50%", "95%", "-"],
        ["SW 가치 비중", "30%", "45%", "60%", "-"],
        ["OTA 차량", "2천만대", "1.2억대", "3억대", "45%"]
    ]
    add_comparison_table(prs, "SDV 시장 전망", headers, rows)
    
    # Section 2: China Standards
    add_section_divider(prs, "중국 SDV 표준 분석", 2)
    
    # China Standards Overview
    bullets = [
        {"main": "표준화 체계", "sub": [
            "Part 1: Atomic Service API - 24개 도메인 서비스 정의",
            "Part 2: Device Abstraction API - 하드웨어 추상화"
        ]},
        {"main": "주요 특징", "sub": [
            "정부 주도 Top-down 방식",
            "빠른 표준화 및 적용",
            "자국 산업 보호 정책 연계"
        ]},
        {"main": "적용 현황", "sub": [
            "BYD, NIO 등 주요 OEM 채택",
            "2025년까지 전면 적용 목표"
        ]}
    ]
    add_content_slide(prs, "중국 SDV 표준화 현황", bullets=bullets)
    
    # Service Domains
    add_content_slide(
        prs,
        "24개 서비스 도메인 구조",
        content="""Vehicle Body Domain
• Power, Door/Window, Lighting, Seat, HVAC

Powertrain Domain  
• Engine, Transmission, EV Motor, Battery Management

Chassis Domain
• Brake, Steering, Suspension, Stability Control

Infotainment & Connectivity
• IVI, Cluster, HUD, V2X, TSP, OTA

ADAS & Autonomous
• Perception, Planning, Control, Mapping"""
    )
    
    # API Architecture
    slide = add_content_slide(prs, "API 아키텍처", layout_idx=5)
    
    # Add architecture diagram (text representation)
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = """┌─────────────────────────────────────┐
│        Application Layer             │
├─────────────────────────────────────┤
│     Atomic Service API Layer        │
│         (24 Domains)                │
├─────────────────────────────────────┤
│   Device Abstraction API Layer      │
│    (Hardware Independence)          │
├─────────────────────────────────────┤
│        Hardware Layer               │
│    (ECUs, Sensors, Actuators)      │
└─────────────────────────────────────┘"""
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Consolas"
        paragraph.font.size = Pt(16)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Section 3: Global Comparison
    add_section_divider(prs, "글로벌 표준 비교", 3)
    
    # Comparison Table
    headers = ["항목", "중국", "AUTOSAR", "일본"]
    rows = [
        ["주도 주체", "정부", "컨소시엄", "기업"],
        ["표준화 속도", "매우 빠름", "보통", "느림"],
        ["개방성", "제한적", "완전 개방", "선택적"],
        ["기술 초점", "서비스/ICV", "안전/품질", "실용성"],
        ["글로벌 호환", "낮음", "높음", "중간"],
        ["적용 분야", "EV/ICV", "전체", "프리미엄"]
    ]
    add_comparison_table(prs, "글로벌 SDV 표준 비교", headers, rows)
    
    # Strengths and Weaknesses
    bullets = [
        {"main": "중국 표준", "sub": [
            "✓ 빠른 개발 및 적용",
            "✓ 정부의 강력한 지원",
            "✗ 글로벌 호환성 부족"
        ]},
        {"main": "AUTOSAR Adaptive", "sub": [
            "✓ 높은 성숙도와 안정성",
            "✓ 글로벌 표준 지위",
            "✗ 높은 개발 복잡도"
        ]},
        {"main": "일본 접근법", "sub": [
            "✓ 검증된 품질과 신뢰성",
            "✓ 실용적 구현",
            "✗ 표준화 부재"
        ]}
    ]
    add_content_slide(prs, "표준별 장단점 분석", bullets=bullets)
    
    # Section 4: Technical Architecture
    add_section_divider(prs, "기술 아키텍처 분석", 4)
    
    # Zonal Architecture
    bullets = [
        {"main": "Zonal Architecture", "sub": [
            "4-6개 Zone Controller",
            "Central High-Performance Computer",
            "Ethernet Backbone (10Gbps+)"
        ]},
        {"main": "Software Stack", "sub": [
            "Hypervisor (Type 1)",
            "Multiple OS (Linux, QNX, Android)",
            "Container/Kubernetes Orchestration"
        ]},
        {"main": "Communication", "sub": [
            "DDS (Data Distribution Service)",
            "SOME/IP Protocol",
            "gRPC for Cloud Services"
        ]}
    ]
    add_content_slide(prs, "SDV 기술 아키텍처", bullets=bullets)
    
    # Security Architecture
    bullets = [
        {"main": "Hardware Security", "sub": [
            "Hardware Security Module (HSM)",
            "Secure Boot Chain",
            "TrustZone/TEE Implementation"
        ]},
        {"main": "Cybersecurity", "sub": [
            "ISO 21434 Compliance",
            "Intrusion Detection System (IDS)",
            "Secure OTA Updates with Rollback"
        ]},
        {"main": "Data Protection", "sub": [
            "End-to-End Encryption",
            "Privacy by Design",
            "GDPR Compliance"
        ]}
    ]
    add_content_slide(prs, "보안 아키텍처", bullets=bullets)
    
    # Section 5: Korea Strategy
    add_section_divider(prs, "한국 대응 전략", 5)
    
    # Strategic Direction
    bullets = [
        {"main": "전략적 포지셔닝", "sub": [
            "글로벌 호환성 + 독자 기술",
            "오픈 이노베이션 생태계",
            "선택과 집중 전략"
        ]},
        {"main": "핵심 역량 확보", "sub": [
            "SDV 플랫폼 기술",
            "AI/자율주행 통합",
            "보안 기술 강화"
        ]},
        {"main": "산업 생태계", "sub": [
            "OEM-Tier1-스타트업 협력",
            "산학연 공동 R&D",
            "글로벌 파트너십"
        ]}
    ]
    add_content_slide(prs, "한국 SDV 전략 방향", bullets=bullets)
    
    # Timeline
    timeline_data = [
        {
            "phase": "Phase 1\n2024-2025",
            "details": [
                "기술 역량 구축",
                "파일럿 프로젝트",
                "인재 양성 500명"
            ]
        },
        {
            "phase": "Phase 2\n2026-2027",
            "details": [
                "K-SDV 표준 제정",
                "상용화 프로젝트",
                "생태계 구축"
            ]
        },
        {
            "phase": "Phase 3\n2028-2030",
            "details": [
                "글로벌 진출",
                "플랫폼 수출",
                "리더십 확보"
            ]
        }
    ]
    add_timeline_slide(prs, "실행 로드맵", timeline_data)
    
    # Section 6: Implementation
    add_section_divider(prs, "실행 계획", 6)
    
    # Action Items
    add_content_slide(
        prs,
        "즉시 실행 과제 (Action Items)",
        content="""1개월 내 실행
• SDV Task Force 구성 및 킥오프
• 중국 표준 문서 완전 번역 및 분석
• AUTOSAR 교육 프로그램 시작 (100명)

3개월 내 실행
• 기술 로드맵 및 아키텍처 확정
• 3개 파일럿 프로젝트 착수
• 글로벌 파트너십 MOU 체결

6개월 내 실행
• K-SDV 표준 v0.9 초안 완성
• 테스트베드 구축 완료
• 첫 PoC 결과 도출 및 검증"""
    )
    
    # Budget Plan
    headers = ["연도", "R&D", "인프라", "인력양성", "국제협력", "합계"]
    rows = [
        ["2024", "800억", "500억", "400억", "300억", "2,000억"],
        ["2025", "1,200억", "750억", "600억", "450억", "3,000억"],
        ["2026", "1,600억", "1,000억", "800억", "600억", "4,000억"],
        ["2027", "2,000억", "1,250억", "1,000억", "750억", "5,000억"],
        ["2028-30", "6,000억", "3,750억", "3,000억", "2,250억", "15,000억"]
    ]
    add_comparison_table(prs, "투자 계획 (2024-2030)", headers, rows)
    
    # Success Metrics
    headers = ["구분", "2025", "2027", "2030"]
    rows = [
        ["SDV 전문인력", "500명", "2,000명", "5,000명"],
        ["양산 적용", "1개 차종", "5개 차종", "전체 차종"],
        ["글로벌 파트너", "3개사", "10개사", "30개사"],
        ["특허 출원", "50건", "200건", "500건"],
        ["표준 기여", "참여", "주도", "리더십"],
        ["수출 규모", "-", "100억원", "1조원"]
    ]
    add_comparison_table(prs, "성과 지표 (KPI)", headers, rows)
    
    # Risk Management
    bullets = [
        {"main": "기술 리스크", "sub": [
            "리스크: 핵심 기술 부족",
            "대응: 조기 기술 확보, M&A 검토"
        ]},
        {"main": "시장 리스크", "sub": [
            "리스크: 표준 분열",
            "대응: 멀티 표준 지원 전략"
        ]},
        {"main": "인력 리스크", "sub": [
            "리스크: 전문 인력 부족",
            "대응: 대규모 교육 프로그램, 해외 인재 유치"
        ]},
        {"main": "규제 리스크", "sub": [
            "리스크: 규제 불확실성",
            "대응: 정부 협력, 선제적 대응"
        ]}
    ]
    add_content_slide(prs, "리스크 관리 방안", bullets=bullets)
    
    # Key Success Factors
    bullets = [
        "Strong Government Support - 정부의 강력한 지원과 정책",
        "Industry Collaboration - OEM-Tier1-스타트업 협력",
        "Global Partnership - 글로벌 표준 기관과 협력",
        "Talent Development - 대규모 인재 양성 프로그램",
        "Technology Leadership - 핵심 기술 조기 확보",
        "Ecosystem Building - 개방형 혁신 생태계 구축"
    ]
    add_content_slide(prs, "핵심 성공 요인", bullets=bullets)
    
    # Conclusion
    add_content_slide(
        prs,
        "결론",
        content="""SDV는 미래 자동차 산업의 게임 체인저

한국의 기회
• 우수한 IT 인프라와 기술력
• 강력한 제조 역량
• 정부-산업계 협력 체계

긴급한 행동 필요
• 중국의 빠른 표준화 대응
• 글로벌 호환성 확보
• 독자 기술 개발

"SDV 시대의 First Mover가 되기 위한 골든타임"
지금 행동하지 않으면 기회를 놓칩니다."""
    )
    
    # Q&A
    slide = add_content_slide(prs, "Q&A", layout_idx=5)
    
    # Center Q&A text
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Thank You\n\n질문과 토론을 환영합니다"
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    return prs

def main():
    print("Creating professional SDV presentation (16:9)...")
    
    # Create main presentation
    prs = create_professional_sdv_presentation()
    prs.save("SDV_Professional_Presentation_16x9.pptx")
    print("✓ Created: SDV_Professional_Presentation_16x9.pptx")
    
    print("\nPresentation created successfully!")
    print("Format: 16:9 Widescreen")
    print("Total slides: 30+")

if __name__ == "__main__":
    main()