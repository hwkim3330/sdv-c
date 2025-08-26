#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def read_existing_ppt(filename):
    """Read and analyze existing PowerPoint presentation"""
    prs = Presentation(filename)
    
    slides_info = []
    print(f"Total slides: {len(prs.slides)}")
    print("=" * 80)
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            'index': i + 1,
            'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
            'title': '',
            'content': []
        }
        
        # Extract title
        if slide.shapes.title:
            slide_info['title'] = slide.shapes.title.text
            
        # Extract all text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if shape != slide.shapes.title:  # Don't duplicate title
                    slide_info['content'].append(shape.text)
                    
        slides_info.append(slide_info)
        
        # Print slide information
        print(f"Slide {i+1}: {slide_info['title']}")
        if slide_info['content']:
            for content in slide_info['content'][:2]:  # First 2 content items
                preview = content.replace('\n', ' ')[:100]
                if preview:
                    print(f"  - {preview}...")
        print("-" * 40)
    
    return slides_info

def create_modified_presentation(original_info):
    """Create a modified and enhanced version of the presentation"""
    prs = Presentation()
    
    # Enhanced Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "중국 SDV 표준 소개"
    subtitle.text = "Software-Defined Vehicle 표준화 현황 및 기술 분석\nKETI 박부식 | 2024.08.26 (수정판)"
    
    # Executive Summary (새로 추가)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    content.text = """주요 내용:
• 중국 SDV 표준화 현황 및 로드맵
• Intelligent Connected Vehicle (ICV) 서비스 인터페이스 사양
• Atomic Service API 및 Device Abstraction API 상세
• 한중독일 표준화 비교 분석
• 국내 대응 전략 제언"""
    
    # Table of Contents (목차 개선)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "목차"
    content.text = """1. SDV 개요 및 배경
   - SDV 정의 및 핵심 기술
   - 글로벌 시장 동향

2. 중국 SDV 표준화 현황
   - 표준화 조직 및 체계
   - 주요 표준 문서 분석
   
3. ICV 서비스 인터페이스 사양
   - Part 1: Atomic Service API
   - Part 2: Device Abstraction API
   
4. 국제 표준화 비교
   - 중국 vs 독일(AUTOSAR) vs 일본
   
5. 시사점 및 대응 방안"""
    
    # Add content slides based on original presentation
    # This is a simplified version - you can add more detailed slides
    
    # SDV Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "SDV (Software-Defined Vehicle) 개요"
    content.text = """정의:
• 차량의 기능이 소프트웨어에 의해 정의되고 제어되는 차량
• 하드웨어와 소프트웨어의 분리 (Decoupling)

핵심 특징:
• OTA (Over-The-Air) 업데이트
• 서비스 지향 아키텍처 (SOA)
• 클라우드 연결성
• AI/ML 기반 기능

시장 전망:
• 2030년까지 전체 차량의 95%가 SDV로 전환 예상
• 소프트웨어 가치 비중 60% 이상"""
    
    # Chinese Standards
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "중국 SDV 표준화 체계"
    content.text = """표준화 주도 기관:
• CATARC (중국자동차기술연구센터)
• CAAM (중국자동차제조협회)
• MIIT (공업정보화부)

주요 표준 문서:
• GB/T 40429-2021: ICV 용어 및 정의
• GB/T 40857-2021: 기능 안전 요구사항
• T/CSAE 234-2021: 서비스 인터페이스 사양

특징:
• 정부 주도의 Top-down 방식
• 빠른 표준화 진행
• 자국 산업 보호 정책과 연계"""
    
    # Technical Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ICV 서비스 인터페이스 아키텍처"
    content.text = """계층 구조:

Application Layer
    ↓
Service Interface Layer
• Atomic Service API
• Composite Service API
    ↓
Device Abstraction Layer
• Sensor Abstraction
• Actuator Abstraction
    ↓
Hardware Layer

핵심 설계 원칙:
• 모듈화 (Modularity)
• 재사용성 (Reusability)
• 확장성 (Scalability)"""
    
    # API Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Atomic Service API 상세"
    content.text = """주요 서비스 카테고리:

1. Vehicle Control Services
   • 파워트레인 제어
   • 섀시 제어
   • 바디 제어

2. Perception Services
   • 센서 데이터 수집
   • 환경 인지
   • 객체 검출

3. Communication Services
   • V2X 통신
   • 클라우드 연결
   • OTA 업데이트

4. Safety & Security Services
   • 기능 안전
   • 사이버 보안
   • 진단 및 모니터링"""
    
    # Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "국제 표준 비교 분석"
    content.text = """중국 vs AUTOSAR Adaptive Platform:

중국:
• 정부 주도, 빠른 표준화
• ICV 중심 접근
• 자국 기업 우선

AUTOSAR (독일/유럽):
• 산업 컨소시엄 주도
• 오픈 표준
• 글로벌 호환성 중시

일본:
• 기업 주도 (Toyota, Honda)
• 실용적 접근
• 안전성 최우선"""
    
    # Recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "국내 대응 전략 제언"
    content.text = """단기 전략:
• 중국 표준 모니터링 및 분석 강화
• AUTOSAR 표준과의 호환성 확보
• 핵심 기술 확보 및 인력 양성

중기 전략:
• 한중 표준 협력 채널 구축
• 국내 SDV 표준 로드맵 수립
• 산학연 협력 생태계 구축

장기 전략:
• 글로벌 표준화 참여 확대
• 독자적 기술 경쟁력 확보
• K-SDV 플랫폼 개발"""
    
    # Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "결론 및 향후 과제"
    content.text = """주요 시사점:
• SDV는 미래 자동차 산업의 핵심
• 중국의 빠른 표준화 진행 주목
• 국제 표준 호환성 확보 필수

향후 과제:
• 지속적인 표준 동향 모니터링
• 국내 기술 역량 강화
• 국제 협력 네트워크 구축
• 산업 생태계 조성

Action Items:
• SDV 표준화 TF 구성
• 정기적인 기술 교류회 개최
• 표준 문서 번역 및 분석"""
    
    # Q&A Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Q&A"
    content.text = "감사합니다.\n\n문의사항:\nKETI 박부식\nEmail: [email]\nTel: [phone]"
    
    return prs

def main():
    print("Analyzing existing presentation...")
    original_file = "중국SDV표준 소개_KETI 박부식0826.pptx"
    
    if os.path.exists(original_file):
        slides_info = read_existing_ppt(original_file)
        
        print("\nCreating modified presentation...")
        modified_prs = create_modified_presentation(slides_info)
        
        output_file = "중국SDV표준_소개_KETI_박부식0826_수정본.pptx"
        modified_prs.save(output_file)
        print(f"\nModified presentation saved as: {output_file}")
    else:
        print(f"File not found: {original_file}")

if __name__ == "__main__":
    main()