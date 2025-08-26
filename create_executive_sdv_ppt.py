#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import datetime

class ProfessionalSDVPresentation:
    def __init__(self):
        self.prs = Presentation()
        self.setup_presentation()
        
        # Color scheme
        self.primary_color = RGBColor(0, 51, 102)      # Dark blue
        self.secondary_color = RGBColor(0, 112, 192)   # Medium blue
        self.accent_color = RGBColor(255, 192, 0)      # Gold
        self.text_color = RGBColor(51, 51, 51)         # Dark gray
        self.light_gray = RGBColor(89, 89, 89)         # Light gray
        
    def setup_presentation(self):
        """Set up 16:9 presentation format"""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def add_cover_slide(self):
        """Add professional cover slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Background gradient effect
        left = top = 0
        width = self.prs.slide_width
        height = self.prs.slide_height
        
        # Add colored rectangle for header
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, width, Inches(1.5)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11.333), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        
        p = title_frame.paragraphs[0]
        p.text = "Software-Defined Vehicle (SDV)"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        p.alignment = PP_ALIGN.CENTER
        
        p = title_frame.add_paragraph()
        p.text = "글로벌 표준화 동향 및 한국 대응 전략"
        p.font.size = Pt(36)
        p.font.color.rgb = self.secondary_color
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            Inches(11.333), Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        
        p = subtitle_frame.paragraphs[0]
        p.text = "중국 SDV 표준 심층 분석 | AUTOSAR 비교 | 실행 로드맵"
        p.font.size = Pt(24)
        p.font.color.rgb = self.light_gray
        p.alignment = PP_ALIGN.CENTER
        
        # Date and version
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11.333), Inches(0.5)
        )
        date_frame = date_box.text_frame
        
        p = date_frame.paragraphs[0]
        p.text = f"{datetime.datetime.now().strftime('%Y년 %m월 %d일')} | Executive Briefing v2.0"
        p.font.size = Pt(18)
        p.font.color.rgb = self.light_gray
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_agenda_slide(self):
        """Add stylish agenda slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Title
        self.add_slide_title(slide, "Agenda")
        
        agenda_items = [
            ("01", "Executive Summary", "핵심 요약 및 주요 메시지"),
            ("02", "SDV 시장 분석", "글로벌 시장 동향 및 전망"),
            ("03", "중국 SDV 표준", "Atomic Service API / Device Abstraction"),
            ("04", "글로벌 표준 비교", "중국 vs AUTOSAR vs 일본"),
            ("05", "기술 아키텍처", "Zonal Architecture / Security"),
            ("06", "한국 대응 전략", "단기-중기-장기 로드맵"),
            ("07", "실행 계획", "Action Items / 투자 계획 / KPI")
        ]
        
        y_start = Inches(1.5)
        for i, (num, title, desc) in enumerate(agenda_items):
            y_pos = y_start + i * Inches(0.8)
            
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1), y_pos,
                Inches(0.6), Inches(0.6)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.secondary_color
            
            text_frame = circle.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = num
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(2), y_pos,
                Inches(4), Inches(0.6)
            )
            text_frame = title_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            
            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(6.5), y_pos,
                Inches(5.5), Inches(0.6)
            )
            text_frame = desc_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.light_gray
            
        return slide
    
    def add_slide_title(self, slide, title_text):
        """Add consistent slide title"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        text_frame = title_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.primary_color
        
        # Add underline
        line = slide.shapes.add_connector(
            1, Inches(0.5), Inches(1.1),
            Inches(12.833), Inches(1.1)
        )
        line.line.color.rgb = self.secondary_color
        line.line.width = Pt(2)
        
    def add_key_message_slide(self):
        """Add key message slide with visual impact"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        self.add_slide_title(slide, "Key Messages")
        
        messages = [
            {
                "icon": "🚗",
                "title": "SDV는 필수",
                "message": "2030년 전체 차량의 95%가 SDV로 전환"
            },
            {
                "icon": "🇨🇳",
                "title": "중국의 도전",
                "message": "정부 주도 빠른 표준화로 시장 선점 시도"
            },
            {
                "icon": "🇰🇷",
                "title": "한국의 기회",
                "message": "글로벌 호환성과 독자 기술의 균형 전략"
            }
        ]
        
        x_start = Inches(1)
        y_pos = Inches(2.5)
        box_width = Inches(3.8)
        
        for i, msg in enumerate(messages):
            x_pos = x_start + i * (box_width + Inches(0.3))
            
            # Message box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y_pos,
                box_width, Inches(3)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            # Icon (using text)
            icon_box = slide.shapes.add_textbox(
                x_pos, y_pos + Inches(0.3),
                box_width, Inches(0.8)
            )
            text_frame = icon_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = msg["icon"]
            p.font.size = Pt(36)
            p.alignment = PP_ALIGN.CENTER
            
            # Title
            title_box = slide.shapes.add_textbox(
                x_pos + Inches(0.2), y_pos + Inches(1.2),
                box_width - Inches(0.4), Inches(0.6)
            )
            text_frame = title_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = msg["title"]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            p.alignment = PP_ALIGN.CENTER
            
            # Message
            msg_box = slide.shapes.add_textbox(
                x_pos + Inches(0.2), y_pos + Inches(1.8),
                box_width - Inches(0.4), Inches(1)
            )
            text_frame = msg_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = msg["message"]
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_color
            p.alignment = PP_ALIGN.CENTER
            
        return slide
    
    def add_market_chart_slide(self):
        """Add market analysis with chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        self.add_slide_title(slide, "SDV 시장 성장 전망")
        
        # Add chart
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(11), Inches(5)
        
        chart_data = CategoryChartData()
        chart_data.categories = ['2024', '2025', '2026', '2027', '2028', '2029', '2030']
        chart_data.add_series('시장 규모 (억 달러)', (650, 980, 1420, 1800, 2200, 2800, 3500))
        chart_data.add_series('SDV 차량 (백만대)', (20, 45, 75, 120, 180, 240, 300))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_title = True
        chart.chart_title.text_frame.text = "SDV 시장 규모 및 차량 대수 전망"
        
        # Add key insights
        insights_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11), Inches(0.8)
        )
        text_frame = insights_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "💡 Key Insight: CAGR 32% 성장 | 2030년 3,500억 달러 규모 | 소프트웨어 가치 비중 60% 도달"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.accent_color
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_architecture_diagram_slide(self):
        """Add technical architecture diagram"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        self.add_slide_title(slide, "SDV 기술 아키텍처")
        
        # Create layered architecture
        layers = [
            ("Applications & Services", self.accent_color, "자율주행, IVI, ADAS, Fleet Management"),
            ("SDV Platform", self.secondary_color, "Service Framework, API Gateway, Data Management"),
            ("Middleware", RGBColor(100, 150, 200), "AUTOSAR Adaptive, DDS, SOME/IP"),
            ("OS & Virtualization", RGBColor(150, 150, 150), "Linux, QNX, Hypervisor, Container"),
            ("Hardware", self.primary_color, "HPC, Zone ECU, Sensors, Actuators")
        ]
        
        y_start = Inches(1.5)
        layer_height = Inches(1.1)
        
        for i, (name, color, desc) in enumerate(layers):
            y_pos = y_start + i * layer_height
            
            # Layer box
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), y_pos,
                Inches(6), layer_height - Inches(0.1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            
            # Layer name
            text_frame = box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(7.5), y_pos,
                Inches(5), layer_height - Inches(0.1)
            )
            text_frame = desc_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_color
            
        return slide
    
    def add_comparison_matrix_slide(self):
        """Add sophisticated comparison matrix"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        self.add_slide_title(slide, "글로벌 SDV 표준 비교 매트릭스")
        
        # Create comparison grid
        data = {
            "headers": ["평가 항목", "중국", "AUTOSAR", "일본"],
            "rows": [
                ["표준화 속도", "●●●●●", "●●●", "●●"],
                ["기술 성숙도", "●●●", "●●●●●", "●●●●"],
                ["글로벌 호환성", "●●", "●●●●●", "●●●"],
                ["생태계 개방성", "●●", "●●●●●", "●●●"],
                ["정부 지원", "●●●●●", "●●", "●●●"],
                ["산업 적용", "●●●●", "●●●", "●●●●●"]
            ]
        }
        
        # Table position
        left = Inches(1)
        top = Inches(1.5)
        cell_width = Inches(2.8)
        cell_height = Inches(0.8)
        
        # Header row
        for i, header in enumerate(data["headers"]):
            x = left + i * cell_width
            
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, top, cell_width - Inches(0.05), cell_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.primary_color
            
            text_frame = box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
        # Data rows
        for row_idx, row_data in enumerate(data["rows"]):
            y = top + (row_idx + 1) * cell_height
            
            for col_idx, cell_text in enumerate(row_data):
                x = left + col_idx * cell_width
                
                box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x, y, cell_width - Inches(0.05), cell_height
                )
                
                # Alternating colors
                if row_idx % 2 == 0:
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGBColor(245, 245, 245)
                else:
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                text_frame = box.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = cell_text
                
                if col_idx == 0:
                    p.font.bold = True
                    p.font.size = Pt(12)
                else:
                    p.font.size = Pt(16)
                    # Color dots based on rating
                    if "●●●●●" in cell_text:
                        p.font.color.rgb = RGBColor(0, 128, 0)
                    elif "●●●●" in cell_text:
                        p.font.color.rgb = RGBColor(100, 150, 0)
                    elif "●●●" in cell_text:
                        p.font.color.rgb = RGBColor(255, 192, 0)
                    else:
                        p.font.color.rgb = RGBColor(255, 100, 0)
                
                p.alignment = PP_ALIGN.CENTER
        
        # Legend
        legend_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5),
            Inches(11), Inches(0.5)
        )
        text_frame = legend_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "● = 낮음  ●●● = 보통  ●●●●● = 높음"
        p.font.size = Pt(12)
        p.font.color.rgb = self.light_gray
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_roadmap_slide(self):
        """Add visual roadmap"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        self.add_slide_title(slide, "한국 SDV 대응 로드맵")
        
        # Timeline arrow
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(1), Inches(3.5),
            Inches(11.5), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.secondary_color
        
        # Phases
        phases = [
            {
                "year": "2024-2025",
                "phase": "Foundation",
                "title": "기반 구축",
                "items": ["기술 역량 확보", "파일럿 프로젝트", "인재 양성"]
            },
            {
                "year": "2026-2027",
                "phase": "Growth",
                "title": "성장 가속",
                "items": ["K-SDV 표준 제정", "상용화 시작", "생태계 구축"]
            },
            {
                "year": "2028-2030",
                "phase": "Leadership",
                "title": "글로벌 리더십",
                "items": ["해외 진출", "플랫폼 수출", "표준 주도"]
            }
        ]
        
        x_start = Inches(1.5)
        for i, phase_data in enumerate(phases):
            x_pos = x_start + i * Inches(3.8)
            
            # Year box
            year_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, Inches(1.8),
                Inches(3.5), Inches(0.6)
            )
            year_box.fill.solid()
            year_box.fill.fore_color.rgb = self.primary_color
            
            text_frame = year_box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = phase_data["year"]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Phase title
            title_box = slide.shapes.add_textbox(
                x_pos, Inches(2.5),
                Inches(3.5), Inches(0.6)
            )
            text_frame = title_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = phase_data["title"]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            p.alignment = PP_ALIGN.CENTER
            
            # Items
            items_box = slide.shapes.add_textbox(
                x_pos, Inches(4.5),
                Inches(3.5), Inches(2)
            )
            text_frame = items_box.text_frame
            text_frame.clear()
            
            for j, item in enumerate(phase_data["items"]):
                if j == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = self.text_color
                p.alignment = PP_ALIGN.LEFT
        
        return slide
    
    def add_action_items_slide(self):
        """Add action items with timeline"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        self.add_slide_title(slide, "Action Items - 즉시 실행 과제")
        
        # Create urgency levels
        urgency_levels = [
            {
                "level": "URGENT",
                "color": RGBColor(255, 0, 0),
                "timeline": "1개월 내",
                "items": [
                    "SDV TF 구성 및 킥오프",
                    "중국 표준 문서 완전 번역",
                    "긴급 인력 양성 계획 수립"
                ]
            },
            {
                "level": "HIGH",
                "color": RGBColor(255, 140, 0),
                "timeline": "3개월 내",
                "items": [
                    "기술 로드맵 확정",
                    "파일럿 프로젝트 3개 착수",
                    "글로벌 파트너십 체결"
                ]
            },
            {
                "level": "MEDIUM",
                "color": RGBColor(255, 192, 0),
                "timeline": "6개월 내",
                "items": [
                    "K-SDV 표준 초안 작성",
                    "테스트베드 구축",
                    "PoC 결과 검증"
                ]
            }
        ]
        
        y_start = Inches(1.5)
        
        for i, level_data in enumerate(urgency_levels):
            y_pos = y_start + i * Inches(1.8)
            
            # Urgency indicator
            indicator = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), y_pos,
                Inches(1.5), Inches(0.5)
            )
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = level_data["color"]
            
            text_frame = indicator.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = level_data["level"]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            # Timeline
            timeline_box = slide.shapes.add_textbox(
                Inches(2.8), y_pos,
                Inches(1.5), Inches(0.5)
            )
            text_frame = timeline_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = level_data["timeline"]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.primary_color
            
            # Items
            items_text = " | ".join(level_data["items"])
            items_box = slide.shapes.add_textbox(
                Inches(4.5), y_pos,
                Inches(7.5), Inches(0.5)
            )
            text_frame = items_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = items_text
            p.font.size = Pt(12)
            p.font.color.rgb = self.text_color
        
        return slide
    
    def add_closing_slide(self):
        """Add impactful closing slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.primary_color
        bg.line.fill.background()
        
        # Main message
        msg_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11.333), Inches(2)
        )
        text_frame = msg_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "The Future is Software-Defined"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = "\n지금이 SDV 시대의 리더가 될 기회입니다"
        p.font.size = Pt(32)
        p.font.color.rgb = self.accent_color
        p.alignment = PP_ALIGN.CENTER
        
        # Call to action
        cta_box = slide.shapes.add_textbox(
            Inches(1), Inches(5),
            Inches(11.333), Inches(1)
        )
        text_frame = cta_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "Act Now or Be Left Behind"
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_presentation(self):
        """Create the complete presentation"""
        # Title and intro
        self.add_cover_slide()
        self.add_agenda_slide()
        self.add_key_message_slide()
        
        # Market analysis
        self.add_market_chart_slide()
        
        # Technical content
        self.add_architecture_diagram_slide()
        self.add_comparison_matrix_slide()
        
        # Strategy and execution
        self.add_roadmap_slide()
        self.add_action_items_slide()
        
        # Closing
        self.add_closing_slide()
        
        return self.prs

def main():
    print("Creating executive SDV presentation...")
    
    presentation = ProfessionalSDVPresentation()
    prs = presentation.create_presentation()
    
    filename = "SDV_Executive_Presentation_Premium.pptx"
    prs.save(filename)
    
    print(f"✅ Successfully created: {filename}")
    print("📊 Format: 16:9 Widescreen (13.333 x 7.5 inches)")
    print("📑 Professional design with visual elements")
    print("🎯 Executive-level content and formatting")

if __name__ == "__main__":
    main()