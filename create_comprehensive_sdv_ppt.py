#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import datetime

def add_slide_with_bullets(prs, title_text, bullets):
    """Helper function to add slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = title_text
    
    # Build bullet text
    bullet_text = ""
    for item in bullets:
        if isinstance(item, str):
            bullet_text += f"• {item}\n"
        elif isinstance(item, dict):
            bullet_text += f"• {item['main']}\n"
            if 'sub' in item:
                for sub in item['sub']:
                    bullet_text += f"  - {sub}\n"
    
    content.text = bullet_text.strip()
    return slide

def create_comprehensive_sdv_presentation():
    """Create a comprehensive SDV presentation combining all insights"""
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SDV (Software-Defined Vehicle)\n완전 분석 보고서"
    subtitle.text = f"중국·독일·일본 표준화 동향 및 한국 대응 전략\n{datetime.datetime.now().strftime('%Y년 %m월 %d일')}"
    
    # 2. Executive Summary
    bullets = [
        "SDV는 미래 자동차 산업의 핵심 패러다임",
        "중국: 정부 주도 빠른 표준화 (ICV 서비스 인터페이스)",
        "독일: AUTOSAR Adaptive Platform (산업 컨소시엄)",
        "일본: 기업 중심 실용적 접근",
        "한국: 글로벌 호환성 + 독자 기술 확보 필요"
    ]
    add_slide_with_bullets(prs, "핵심 요약", bullets)
    
    # 3. Agenda
    bullets = [
        {"main": "Part 1: SDV 개념 및 시장 전망", "sub": ["SDV 정의", "핵심 기술", "시장 규모"]},
        {"main": "Part 2: 글로벌 표준화 현황", "sub": ["중국", "독일/유럽", "일본", "미국"]},
        {"main": "Part 3: 중국 SDV 표준 심층 분석", "sub": ["Atomic Service API", "Device Abstraction API"]},
        {"main": "Part 4: 기술 비교 분석", "sub": ["표준별 장단점", "호환성 이슈"]},
        {"main": "Part 5: 한국 대응 전략", "sub": ["단기/중기/장기 로드맵", "Action Items"]}
    ]
    add_slide_with_bullets(prs, "목차", bullets)
    
    # 4. PART 1 - SDV Overview
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "PART 1\nSDV 개념 및 시장 전망"
    
    # 5. SDV Definition
    bullets = [
        "Software-Defined Vehicle = 소프트웨어가 차량의 기능과 성능을 정의",
        {"main": "핵심 특징", "sub": [
            "하드웨어와 소프트웨어 분리 (Decoupling)",
            "OTA(Over-The-Air) 업데이트",
            "서비스 지향 아키텍처 (SOA)",
            "클라우드 네이티브 기술 적용"
        ]},
        {"main": "기존 차량 vs SDV", "sub": [
            "기존: ECU별 고정 기능, 하드웨어 종속",
            "SDV: 중앙 집중형, 소프트웨어 정의, 유연한 기능 추가"
        ]}
    ]
    add_slide_with_bullets(prs, "SDV 정의 및 특징", bullets)
    
    # 6. Market Outlook
    bullets = [
        "2024년: 전체 차량의 15% SDV 탑재",
        "2027년: 50% 이상 SDV 전환 예상",
        "2030년: 95% SDV 시대 도래",
        {"main": "시장 규모", "sub": [
            "2024년: 650억 달러",
            "2030년: 3,500억 달러 (CAGR 32%)"
        ]},
        {"main": "주요 플레이어", "sub": [
            "OEM: Tesla, BYD, Volkswagen, Hyundai",
            "Tech: Google, Apple, Nvidia, Qualcomm"
        ]}
    ]
    add_slide_with_bullets(prs, "SDV 시장 전망", bullets)
    
    # 7. PART 2 - Global Standards
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "PART 2\n글로벌 표준화 현황"
    
    # 8. China Standards
    bullets = [
        {"main": "표준화 주체", "sub": [
            "MIIT (공업정보화부)",
            "CATARC (중국자동차기술연구센터)",
            "CAAM (중국자동차제조협회)"
        ]},
        {"main": "주요 표준", "sub": [
            "GB/T 40429-2021: ICV 용어 및 정의",
            "GB/T 40857-2021: 기능 안전 요구사항",
            "T/CSAE 234-2021: 서비스 인터페이스 사양"
        ]},
        {"main": "특징", "sub": [
            "정부 주도 Top-down 방식",
            "빠른 표준화 진행",
            "자국 산업 보호와 연계"
        ]}
    ]
    add_slide_with_bullets(prs, "중국 SDV 표준화", bullets)
    
    # 9. Germany/Europe Standards
    bullets = [
        {"main": "AUTOSAR Adaptive Platform", "sub": [
            "2017년 첫 릴리즈",
            "POSIX 기반 OS",
            "Service-Oriented Architecture"
        ]},
        {"main": "주요 특징", "sub": [
            "산업 컨소시엄 주도 (140+ 멤버)",
            "오픈 표준",
            "기능 안전 (ISO 26262) 중심"
        ]},
        {"main": "적용 사례", "sub": [
            "BMW: iDrive 8/9",
            "Mercedes: MB.OS",
            "Volkswagen: VW.OS"
        ]}
    ]
    add_slide_with_bullets(prs, "독일/유럽 SDV 표준화", bullets)
    
    # 10. Japan Standards
    bullets = [
        {"main": "접근 방식", "sub": [
            "기업 주도 (Toyota, Honda, Nissan)",
            "실용성 중심",
            "안전 최우선"
        ]},
        {"main": "주요 이니셔티브", "sub": [
            "Arene OS (Toyota)",
            "Honda CONNECT 3.0",
            "Nissan Intelligent Mobility"
        ]},
        {"main": "특징", "sub": [
            "점진적 전환 전략",
            "하이브리드 아키텍처",
            "품질/신뢰성 중시"
        ]}
    ]
    add_slide_with_bullets(prs, "일본 SDV 표준화", bullets)
    
    # 11. PART 3 - China Deep Dive
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "PART 3\n중국 SDV 표준 심층 분석"
    
    # 12. Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "중국 SDV 아키텍처 구조"
    content.text = """Application Layer
├── Smart Apps / Services
│
Service Interface Layer
├── Atomic Service API (Part 1)
│   ├── 24 Service Domains
│   └── Standardized Methods & Events
│
├── Device Abstraction API (Part 2)
│   ├── Sensor Abstraction
│   └── Actuator Abstraction
│
Hardware Layer
├── ECUs / Sensors / Actuators
└── Vehicle Network (CAN/Ethernet)"""
    
    # 13. Atomic Service API Details
    bullets = [
        {"main": "24개 서비스 도메인", "sub": [
            "Vehicle Body: Power, Door, Lighting, Seat",
            "Powertrain: Engine, Transmission, EV Motor",
            "Chassis: Brake, Steering, Suspension",
            "Infotainment: IVI, Cluster, HUD",
            "ADAS: Perception, Planning, Control",
            "Connectivity: V2X, TSP, OTA"
        ]},
        {"main": "API 설계 원칙", "sub": [
            "RESTful 스타일",
            "JSON 메시지 포맷",
            "비동기 이벤트 지원"
        ]}
    ]
    add_slide_with_bullets(prs, "Atomic Service API 상세", bullets)
    
    # 14. API Example
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "API 호출 예시"
    content.text = """Power System Service:
{
  "serviceId": "power.system",
  "method": "setPowerMode",
  "params": {
    "mode": "COMFORT",
    "profile": "ECO"
  }
}

Response:
{
  "result": "SUCCESS",
  "data": {
    "actualMode": "COMFORT",
    "batteryLevel": 78,
    "estimatedRange": 420
  }
}"""
    
    # 15. Device Abstraction Details
    bullets = [
        {"main": "추상화 계층", "sub": [
            "Physical Device → Abstract Interface → Service",
            "하드웨어 독립성 확보",
            "벤더 중립적 인터페이스"
        ]},
        {"main": "주요 인터페이스", "sub": [
            "Sensor: read(), calibrate(), getStatus()",
            "Actuator: write(), control(), feedback()",
            "ECU: configure(), diagnose(), update()"
        ]},
        {"main": "이점", "sub": [
            "부품 교체 용이",
            "멀티 벤더 지원",
            "신속한 프로토타이핑"
        ]}
    ]
    add_slide_with_bullets(prs, "Device Abstraction API 상세", bullets)
    
    # 16. PART 4 - Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "PART 4\n기술 비교 분석"
    
    # 17. Standards Comparison Table
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "표준 비교 매트릭스"
    content.text = """항목          중국         AUTOSAR      일본
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
주도          정부         컨소시엄      기업
개방성        제한적       완전개방      선택적
표준화속도    매우빠름     보통         느림
기술초점      서비스       안전         품질
생태계        폐쇄적       개방적       하이브리드
호환성        낮음         높음         중간
적용분야      ICV/EV       전체         프리미엄"""
    
    # 18. Strengths and Weaknesses
    bullets = [
        {"main": "중국 표준", "sub": [
            "강점: 빠른 개발, 정부 지원, 대규모 시장",
            "약점: 글로벌 호환성, 폐쇄적 생태계"
        ]},
        {"main": "AUTOSAR", "sub": [
            "강점: 성숙도, 글로벌 표준, 안전성",
            "약점: 복잡성, 높은 진입장벽"
        ]},
        {"main": "일본 방식", "sub": [
            "강점: 품질, 신뢰성, 실용성",
            "약점: 표준화 부재, 확장성 제한"
        ]}
    ]
    add_slide_with_bullets(prs, "표준별 강점과 약점", bullets)
    
    # 19. PART 5 - Korea Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "PART 5\n한국 대응 전략"
    
    # 20. Short-term Strategy (2024-2025)
    bullets = [
        {"main": "기술 확보", "sub": [
            "AUTOSAR Adaptive Platform 역량 구축",
            "중국 표준 분석 및 매핑 가이드 개발",
            "핵심 인재 양성 (연 500명)"
        ]},
        {"main": "파일럿 프로젝트", "sub": [
            "K-SDV 참조 모델 개발",
            "OEM-Tier1 공동 PoC",
            "테스트베드 구축"
        ]},
        {"main": "국제 협력", "sub": [
            "AUTOSAR 멤버십 확대",
            "중국 표준 기관과 MOU",
            "글로벌 기업 파트너십"
        ]}
    ]
    add_slide_with_bullets(prs, "단기 전략 (2024-2025)", bullets)
    
    # 21. Mid-term Strategy (2026-2027)
    bullets = [
        {"main": "표준 개발", "sub": [
            "K-SDV 표준 v1.0 제정",
            "AUTOSAR 프로파일 정의",
            "중국 표준 브리지 개발"
        ]},
        {"main": "생태계 구축", "sub": [
            "SDV 개발 플랫폼 구축",
            "오픈소스 프로젝트 런칭",
            "스타트업 인큐베이션"
        ]},
        {"main": "상용화", "sub": [
            "2-3개 양산 프로젝트",
            "After-market 솔루션",
            "수출 모델 개발"
        ]}
    ]
    add_slide_with_bullets(prs, "중기 전략 (2026-2027)", bullets)
    
    # 22. Long-term Vision (2028-2030)
    bullets = [
        {"main": "글로벌 리더십", "sub": [
            "K-SDV 국제 표준화",
            "독자 기술 IP 확보",
            "글로벌 Top 3 진입"
        ]},
        {"main": "차세대 기술", "sub": [
            "AI 기반 자율 최적화",
            "양자 암호 보안",
            "6G 연계 서비스"
        ]},
        {"main": "새로운 비즈니스", "sub": [
            "SDV 플랫폼 수출",
            "소프트웨어 라이선스",
            "데이터 서비스"
        ]}
    ]
    add_slide_with_bullets(prs, "장기 비전 (2028-2030)", bullets)
    
    # 23. Action Items
    bullets = [
        "즉시 실행 (1개월 내)",
        "• SDV TF 구성 및 킥오프",
        "• 중국 표준 문서 번역 완료",
        "• AUTOSAR 교육 프로그램 시작",
        "",
        "단기 실행 (3개월 내)",
        "• 기술 로드맵 수립",
        "• 파일럿 프로젝트 선정",
        "• 국제 협력 채널 구축",
        "",
        "중기 실행 (6개월 내)",
        "• K-SDV 표준 초안 작성",
        "• 테스트베드 구축 완료",
        "• 첫 PoC 결과 도출"
    ]
    add_slide_with_bullets(prs, "Action Items", bullets)
    
    # 24. Risk Management
    bullets = [
        {"main": "기술 리스크", "sub": [
            "대응: 핵심 기술 조기 확보",
            "대응: 다중 기술 트랙 운영"
        ]},
        {"main": "시장 리스크", "sub": [
            "대응: 유연한 표준 전략",
            "대응: 글로벌 파트너십"
        ]},
        {"main": "인력 리스크", "sub": [
            "대응: 체계적 교육 프로그램",
            "대응: 해외 인재 유치"
        ]},
        {"main": "규제 리스크", "sub": [
            "대응: 정부 협력 강화",
            "대응: 선제적 규제 대응"
        ]}
    ]
    add_slide_with_bullets(prs, "리스크 관리", bullets)
    
    # 25. Budget Estimation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "예상 투자 규모"
    content.text = """2024-2030 총 투자 규모: 3조원

연도별 투자 계획:
• 2024: 2,000억원 (기반 구축)
• 2025: 3,000억원 (기술 개발)
• 2026: 4,000억원 (상용화)
• 2027: 5,000억원 (확산)
• 2028-2030: 연 5,000억원 (글로벌화)

투자 분야:
• R&D: 40%
• 인프라: 25%
• 인력 양성: 20%
• 국제 협력: 15%"""
    
    # 26. Success Metrics
    bullets = [
        {"main": "2025년 목표", "sub": [
            "SDV 전문인력 500명",
            "파일럿 프로젝트 3개",
            "특허 출원 50건"
        ]},
        {"main": "2027년 목표", "sub": [
            "양산 적용 5개 차종",
            "글로벌 파트너 10개사",
            "수출 100억원"
        ]},
        {"main": "2030년 목표", "sub": [
            "글로벌 시장 점유율 15%",
            "SDV 플랫폼 수출 10개국",
            "매출 1조원"
        ]}
    ]
    add_slide_with_bullets(prs, "성공 지표", bullets)
    
    # 27. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "결론"
    content.text = """SDV는 선택이 아닌 필수

핵심 성공 요인:
• 빠른 실행
• 글로벌 협력
• 독자 기술 확보
• 생태계 구축

한국의 기회:
• 강력한 IT 인프라
• 우수한 인재풀
• 완성차-부품사 협력 문화
• 정부의 적극적 지원

"First Mover가 되기 위한 골든타임"
지금이 행동할 때입니다."""
    
    # 28. Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Q&A"
    content.text = "감사합니다.\n\n질문과 토론을 환영합니다."
    
    return prs

def create_technical_deep_dive():
    """Create technical deep dive presentation"""
    prs = Presentation()
    
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SDV 기술 심화 분석"
    subtitle.text = "Architecture, API, Implementation Details"
    
    # Technical Architecture
    bullets = [
        {"main": "Zonal Architecture", "sub": [
            "4-6 Zone Controllers",
            "Central Computing Unit",
            "Ethernet Backbone (TSN)"
        ]},
        {"main": "Software Stack", "sub": [
            "Hypervisor (Type 1)",
            "Multiple OS (Linux, RTOS, Android)",
            "Container/Kubernetes"
        ]},
        {"main": "Communication", "sub": [
            "DDS (Data Distribution Service)",
            "SOME/IP",
            "gRPC/REST"
        ]}
    ]
    add_slide_with_bullets(prs, "Technical Architecture", bullets)
    
    # API Design Patterns
    bullets = [
        "Service Mesh Architecture",
        "API Gateway Pattern",
        "Event-Driven Architecture",
        "Microservices Design",
        "Circuit Breaker Pattern",
        "Saga Pattern for Transactions"
    ]
    add_slide_with_bullets(prs, "API Design Patterns", bullets)
    
    # Security Architecture
    bullets = [
        {"main": "Security Layers", "sub": [
            "Hardware Security Module (HSM)",
            "Secure Boot Chain",
            "TrustZone/TEE"
        ]},
        {"main": "Cybersecurity", "sub": [
            "ISO 21434 Compliance",
            "Intrusion Detection System",
            "Secure OTA Updates"
        ]}
    ]
    add_slide_with_bullets(prs, "Security Architecture", bullets)
    
    # Save presentation
    return prs

def main():
    print("Creating comprehensive SDV presentations...")
    
    # Create main comprehensive presentation
    comprehensive_prs = create_comprehensive_sdv_presentation()
    comprehensive_prs.save("SDV_종합분석_보고서.pptx")
    print("Created: SDV_종합분석_보고서.pptx")
    
    # Create technical deep dive
    technical_prs = create_technical_deep_dive()
    technical_prs.save("SDV_기술심화_분석.pptx")
    print("Created: SDV_기술심화_분석.pptx")
    
    print("\nAll presentations created successfully!")

if __name__ == "__main__":
    main()