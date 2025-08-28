#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from datetime import datetime

def add_table_slide(prs, title, table_data, headers):
    """테이블 슬라이드 추가"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Table
    rows = len(table_data) + 1  # +1 for header
    cols = len(headers)
    left = Inches(1)
    top = Inches(2)
    width = Inches(14)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(14 / cols)
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 84, 159)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add data
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.LEFT
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    return slide

def add_chart_slide(prs, title, chart_type, categories, series_data):
    """차트 슬라이드 추가"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)
    
    x, y, cx, cy = Inches(1), Inches(2), Inches(14), Inches(6)
    chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
    
    # Chart styling
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    return slide

def add_diagram_slide(prs, title, diagram_type="architecture"):
    """다이어그램 슬라이드 추가"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    if diagram_type == "architecture":
        # SDV Architecture Layers
        layers = [
            ("Application Layer", "사용자 앱, 인포테인먼트, ADAS", RGBColor(0, 176, 240)),
            ("Service Layer", "API Gateway, 마이크로서비스, 데이터 관리", RGBColor(0, 112, 192)),
            ("Middleware Layer", "AUTOSAR AP, ROS, DDS", RGBColor(0, 84, 159)),
            ("OS/Hypervisor", "Linux, QNX, 하이퍼바이저", RGBColor(75, 172, 198)),
            ("Hardware Layer", "ECU, 센서, 액추에이터", RGBColor(128, 128, 128))
        ]
        
        y_start = 2
        for i, (layer_name, description, color) in enumerate(layers):
            # Layer box
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(2), Inches(y_start + i * 1.1),
                Inches(12), Inches(1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = RGBColor(255, 255, 255)
            box.line.width = Pt(2)
            
            # Layer text
            text_frame = box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = f"{layer_name}: {description}"
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            
    elif diagram_type == "ecosystem":
        # SDV Ecosystem
        center_x, center_y = 8, 4.5
        radius = 2.5
        
        # Central circle - SDV Platform
        center = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_x - 1.5), Inches(center_y - 1.5),
            Inches(3), Inches(3)
        )
        center.fill.solid()
        center.fill.fore_color.rgb = RGBColor(0, 84, 159)
        text_frame = center.text_frame
        text_frame.text = "SDV\n플랫폼"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Surrounding elements
        elements = [
            ("OEM", 0, RGBColor(255, 192, 0)),
            ("Tier 1", 60, RGBColor(112, 173, 71)),
            ("SW 벤더", 120, RGBColor(91, 155, 213)),
            ("클라우드", 180, RGBColor(255, 0, 0)),
            ("반도체", 240, RGBColor(112, 48, 160)),
            ("서비스", 300, RGBColor(237, 125, 49))
        ]
        
        import math
        for name, angle, color in elements:
            # Calculate position
            rad = math.radians(angle)
            x = center_x + radius * math.cos(rad) - 0.75
            y = center_y + radius * math.sin(rad) - 0.75
            
            # Add circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(1.5), Inches(1.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            
            # Add text
            text_frame = circle.text_frame
            text_frame.text = name
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add connector line
            connector = slide.shapes.add_connector(
                1,  # Straight connector
                Inches(center_x), Inches(center_y),
                Inches(x + 0.75), Inches(y + 0.75)
            )
            connector.line.color.rgb = RGBColor(200, 200, 200)
            connector.line.width = Pt(1)
    
    elif diagram_type == "timeline":
        # Timeline
        timeline_y = 5
        start_x = 2
        end_x = 14
        
        # Main timeline
        line = slide.shapes.add_connector(
            1,
            Inches(start_x), Inches(timeline_y),
            Inches(end_x), Inches(timeline_y)
        )
        line.line.color.rgb = RGBColor(0, 84, 159)
        line.line.width = Pt(3)
        
        # Timeline points
        milestones = [
            (2025, "기반 구축", "표준화 착수"),
            (2027, "확산 단계", "상용화 시작"),
            (2029, "성숙 단계", "완전 자율주행"),
            (2030, "글로벌 리더", "시장 주도")
        ]
        
        for i, (year, title, desc) in enumerate(milestones):
            x_pos = start_x + (end_x - start_x) * i / (len(milestones) - 1)
            
            # Milestone circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos - 0.2), Inches(timeline_y - 0.2),
                Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(255, 192, 0)
            
            # Year text above
            year_box = slide.shapes.add_textbox(
                Inches(x_pos - 0.5), Inches(timeline_y - 1.2),
                Inches(1), Inches(0.5)
            )
            year_box.text_frame.text = str(year)
            year_box.text_frame.paragraphs[0].font.bold = True
            year_box.text_frame.paragraphs[0].font.size = Pt(14)
            year_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Description below
            desc_box = slide.shapes.add_textbox(
                Inches(x_pos - 1.5), Inches(timeline_y + 0.5),
                Inches(3), Inches(1)
            )
            desc_box.text_frame.text = f"{title}\n{desc}"
            desc_box.text_frame.paragraphs[0].font.size = Pt(11)
            desc_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    elif diagram_type == "process":
        # Process Flow
        processes = [
            ("요구사항\n분석", RGBColor(0, 176, 240)),
            ("아키텍처\n설계", RGBColor(0, 112, 192)),
            ("개발 및\n통합", RGBColor(0, 84, 159)),
            ("검증 및\n테스트", RGBColor(112, 173, 71)),
            ("배포 및\n운영", RGBColor(255, 192, 0))
        ]
        
        x_start = 1.5
        y_pos = 4
        box_width = 2.5
        spacing = 0.3
        
        for i, (process_name, color) in enumerate(processes):
            x_pos = x_start + i * (box_width + spacing)
            
            # Process box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(y_pos),
                Inches(box_width), Inches(1.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            
            # Process text
            text_frame = box.text_frame
            text_frame.text = process_name
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Arrow to next
            if i < len(processes) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(x_pos + box_width), Inches(y_pos + 0.5),
                    Inches(spacing), Inches(0.5)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
    
    return slide

def create_advanced_sdv_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Slide 1: Title Slide with Design
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Software-Defined Vehicle (SDV)"
    subtitle.text = "글로벌 표준화 동향 및 한국 대응 전략\n차세대 모빌리티 혁신을 위한 종합 로드맵\n2025년 1월"
    
    # Add decorative shapes to title slide
    rect1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(8.5),
        Inches(16), Inches(0.5)
    )
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = RGBColor(0, 84, 159)
    rect1.line.fill.background()
    
    # Slide 2: Executive Summary with Icons
    bullet_slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Executive Summary"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Key Points with boxes
    key_points = [
        ("시장 규모", "2030년 1,500억 달러", "연평균 성장률 25%", RGBColor(0, 176, 240)),
        ("기술 전환", "SW 가치 비중 60%", "OTA 업데이트 표준화", RGBColor(112, 173, 71)),
        ("경쟁 현황", "중국 표준 선점", "독일 AUTOSAR 주도", RGBColor(255, 192, 0)),
        ("한국 전략", "K-SDV 플랫폼", "5년간 6,500억 투자", RGBColor(237, 125, 49))
    ]
    
    x_start = 1
    y_start = 2.5
    for i, (title, main, sub, color) in enumerate(key_points):
        x = x_start + (i % 2) * 7.5
        y = y_start + (i // 2) * 3
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(6.5), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        # Text
        text_frame = box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        p = text_frame.add_paragraph()
        p.text = main
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        p = text_frame.add_paragraph()
        p.text = sub
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Slide 3: SDV Architecture Diagram
    add_diagram_slide(prs, "SDV 기술 아키텍처", "architecture")
    
    # Slide 4: Market Growth Chart
    categories = ['2024', '2025', '2026', '2027', '2028', '2029', '2030']
    series_data = {
        '글로벌 시장': [250, 350, 500, 700, 950, 1200, 1500],
        '한국 시장': [15, 25, 40, 65, 95, 130, 180],
        '중국 시장': [80, 120, 180, 260, 350, 450, 550]
    }
    add_chart_slide(prs, "SDV 시장 성장 전망 (억 달러)", XL_CHART_TYPE.COLUMN_CLUSTERED, categories, series_data)
    
    # Slide 5: Global Standards Comparison Table
    headers = ['국가', '표준 명칭', '주요 특징', '완성도', '채택률']
    table_data = [
        ['중국', 'C-ICVS V4', '정부 주도 통합 표준', '90%', '높음'],
        ['독일', 'AUTOSAR AP', '업계 컨소시엄', '95%', '매우 높음'],
        ['일본', 'Vehicle OS', '안전성 중심', '85%', '중간'],
        ['미국', 'SDV Alliance', '오픈소스 기반', '75%', '증가중'],
        ['한국', 'K-SDV', '개발 중', '30%', '초기']
    ]
    add_table_slide(prs, "글로벌 SDV 표준화 현황 비교", table_data, headers)
    
    # Slide 6: SDV Ecosystem Diagram
    add_diagram_slide(prs, "SDV 생태계 구조", "ecosystem")
    
    # Slide 7: Technology Stack with SmartArt style
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SDV 핵심 기술 스택"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Create pyramid-like structure
    levels = [
        ("클라우드 서비스", ["AWS", "Azure", "GCP", "알리바바"], RGBColor(0, 176, 240)),
        ("통신 프로토콜", ["5G/6G", "V2X", "WiFi 6", "Ethernet"], RGBColor(0, 140, 200)),
        ("미들웨어", ["SOME/IP", "DDS", "ROS2", "MQTT"], RGBColor(0, 112, 192)),
        ("운영체제", ["Linux", "QNX", "Android Auto", "AGL"], RGBColor(0, 84, 159)),
        ("하드웨어", ["GPU", "NPU", "MCU", "센서"], RGBColor(75, 172, 198))
    ]
    
    y_pos = 2
    for i, (category, items, color) in enumerate(levels):
        width = 10 + i * 1.5  # Pyramid effect
        x_pos = (16 - width) / 2
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            Inches(x_pos), Inches(y_pos + i * 1.1),
            Inches(width), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        text_frame = box.text_frame
        text_frame.text = f"{category}: {', '.join(items)}"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 8: Investment Distribution Pie Chart
    categories = ['플랫폼 개발', '인프라 구축', '인력 양성', '표준화', 'R&D']
    series_data = {
        '투자 비중': [45, 25, 15, 8, 7]
    }
    add_chart_slide(prs, "SDV 투자 배분 계획 (%)", XL_CHART_TYPE.PIE, categories, series_data)
    
    # Slide 9: Timeline Roadmap
    add_diagram_slide(prs, "SDV 구현 로드맵", "timeline")
    
    # Slide 10: SWOT Analysis with Matrix
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "한국 SDV 산업 SWOT 분석"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # SWOT Matrix
    swot_data = [
        ("강점 (S)", ["우수한 IT 인프라", "5G 네트워크 선도", "반도체 기술력", "제조업 경쟁력"], 
         RGBColor(112, 173, 71)),
        ("약점 (W)", ["SW 플랫폼 부족", "표준 주도권 미흡", "전문 인력 부족", "생태계 미성숙"], 
         RGBColor(255, 192, 0)),
        ("기회 (O)", ["시장 급성장", "정부 지원 확대", "신기술 융합", "글로벌 협력"], 
         RGBColor(91, 155, 213)),
        ("위협 (T)", ["중국 추격", "기술 격차", "규제 불확실성", "투자 리스크"], 
         RGBColor(237, 125, 49))
    ]
    
    for i, (title, items, color) in enumerate(swot_data):
        x = 1 + (i % 2) * 7.5
        y = 2 + (i // 2) * 3.5
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(6.5), Inches(3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        text_frame = box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.level = 0
    
    # Slide 11: Process Flow
    add_diagram_slide(prs, "SDV 개발 프로세스", "process")
    
    # Slide 12: Performance Metrics Dashboard
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SDV 성과 지표 대시보드"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # KPI Cards
    kpis = [
        ("OTA 성공률", "99.9%", "▲ 2.3%", RGBColor(112, 173, 71)),
        ("서비스 가용성", "99.95%", "▲ 0.5%", RGBColor(112, 173, 71)),
        ("보안 사고", "0건", "─ 0", RGBColor(91, 155, 213)),
        ("고객 만족도", "4.5/5.0", "▲ 0.3", RGBColor(112, 173, 71)),
        ("시장 점유율", "15%", "▲ 3%", RGBColor(112, 173, 71)),
        ("개발 속도", "2주 스프린트", "▼ 3일", RGBColor(255, 192, 0))
    ]
    
    for i, (metric, value, change, color) in enumerate(kpis):
        x = 1.5 + (i % 3) * 5
        y = 2 + (i // 3) * 3
        
        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(4), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        
        # Metric name
        metric_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.2),
            Inches(3.6), Inches(0.5)
        )
        metric_box.text_frame.text = metric
        metric_box.text_frame.paragraphs[0].font.size = Pt(14)
        metric_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        # Value
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.7),
            Inches(3.6), Inches(1)
        )
        value_box.text_frame.text = value
        value_box.text_frame.paragraphs[0].font.size = Pt(28)
        value_box.text_frame.paragraphs[0].font.bold = True
        value_box.text_frame.paragraphs[0].font.color.rgb = color
        
        # Change indicator
        change_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 1.7),
            Inches(3.6), Inches(0.5)
        )
        change_box.text_frame.text = change
        change_box.text_frame.paragraphs[0].font.size = Pt(12)
        change_box.text_frame.paragraphs[0].font.color.rgb = color
    
    # Slide 13: Competitive Analysis Radar Chart
    categories = ['2025', '2026', '2027', '2028', '2029', '2030']
    series_data = {
        '한국': [30, 45, 60, 75, 85, 95],
        '중국': [70, 75, 80, 85, 90, 95],
        '독일': [80, 85, 88, 90, 92, 95],
        '일본': [60, 65, 70, 75, 80, 85]
    }
    add_chart_slide(prs, "국가별 SDV 기술 성숙도 전망 (%)", XL_CHART_TYPE.LINE, categories, series_data)
    
    # Slide 14: Security Architecture
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SDV 보안 아키텍처"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Security Layers (Onion model)
    center_x, center_y = 8, 4.5
    layers = [
        ("하드웨어 보안", 4, RGBColor(200, 200, 200)),
        ("OS 보안", 3.2, RGBColor(150, 150, 150)),
        ("네트워크 보안", 2.4, RGBColor(100, 100, 100)),
        ("애플리케이션 보안", 1.6, RGBColor(50, 50, 50)),
        ("데이터 보안", 0.8, RGBColor(0, 0, 0))
    ]
    
    for name, radius, color in layers:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(center_x - radius), Inches(center_y - radius),
            Inches(radius * 2), Inches(radius * 2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.fill.transparency = 0.3
        
        # Add label
        label_box = slide.shapes.add_textbox(
            Inches(center_x + radius + 0.2), Inches(center_y - 0.3),
            Inches(2), Inches(0.5)
        )
        label_box.text_frame.text = name
        label_box.text_frame.paragraphs[0].font.size = Pt(11)
        label_box.text_frame.paragraphs[0].font.bold = True
    
    # Add security measures around
    measures = [
        ("암호화", 2, 2, RGBColor(255, 0, 0)),
        ("인증", 14, 2, RGBColor(0, 255, 0)),
        ("침입탐지", 2, 7, RGBColor(0, 0, 255)),
        ("접근제어", 14, 7, RGBColor(255, 255, 0))
    ]
    
    for measure, x, y, color in measures:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(1.5), Inches(0.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.fill.transparency = 0.7
        
        text_frame = box.text_frame
        text_frame.text = measure
        text_frame.paragraphs[0].font.size = Pt(10)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 15: API Standards Comparison
    headers = ['API 유형', '중국 표준', '독일 표준', '한국 목표']
    table_data = [
        ['원자 서비스 API', 'V4 Beta 완성', 'AUTOSAR AP', '개발 중'],
        ['디바이스 추상화', '표준화 완료', '부분 표준화', '설계 단계'],
        ['데이터 관리', 'DDS 기반', 'SOME/IP', '하이브리드'],
        ['보안 프로토콜', 'TLS 1.3', 'SecOC', 'TLS + SecOC'],
        ['실시간 통신', '5ms 이하', '1ms 이하', '1ms 목표']
    ]
    add_table_slide(prs, "API 표준 비교 분석", table_data, headers)
    
    # Slide 16: Business Model Innovation
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SDV 비즈니스 모델 혁신"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Revenue streams
    streams = [
        ("하드웨어 판매", "30%", "기존 모델", RGBColor(128, 128, 128)),
        ("SW 라이선스", "25%", "신규 수익", RGBColor(0, 176, 240)),
        ("구독 서비스", "20%", "지속 수익", RGBColor(112, 173, 71)),
        ("데이터 서비스", "15%", "미래 성장", RGBColor(255, 192, 0)),
        ("앱 마켓", "10%", "생태계 수익", RGBColor(237, 125, 49))
    ]
    
    # Create flow diagram
    y_start = 2.5
    for i, (stream, percent, type, color) in enumerate(streams):
        # Stream box
        box = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(2), Inches(y_start + i * 1.2),
            Inches(4), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        # Stream text
        text_frame = box.text_frame
        text_frame.text = stream
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Percentage
        pct_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(6.5), Inches(y_start + i * 1.2),
            Inches(1), Inches(1)
        )
        pct_box.fill.solid()
        pct_box.fill.fore_color.rgb = color
        
        text_frame = pct_box.text_frame
        text_frame.text = percent
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Type description
        type_box = slide.shapes.add_textbox(
            Inches(8), Inches(y_start + i * 1.2 + 0.25),
            Inches(3), Inches(0.5)
        )
        type_box.text_frame.text = type
        type_box.text_frame.paragraphs[0].font.size = Pt(11)
        type_box.text_frame.paragraphs[0].font.color.rgb = color
    
    # Total revenue indicator
    total_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.5), Inches(3),
        Inches(3.5), Inches(3)
    )
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = RGBColor(0, 84, 159)
    
    text_frame = total_box.text_frame
    text_frame.text = "총 수익\n목표\n\n연 10조원\n(2030년)"
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 17: Partnership Strategy Matrix
    headers = ['파트너 유형', '핵심 역할', '협력 모델', '우선순위']
    table_data = [
        ['글로벌 OEM', '공동 플랫폼 개발', 'JV/전략적 제휴', '★★★★★'],
        ['빅테크 기업', '클라우드/AI 기술', '기술 라이선싱', '★★★★★'],
        ['Tier 1 공급사', '부품 통합', '수직계열화', '★★★★'],
        ['스타트업', '혁신 기술', 'M&A/투자', '★★★'],
        ['대학/연구소', 'R&D 협력', '산학협력', '★★★']
    ]
    add_table_slide(prs, "전략적 파트너십 구축 계획", table_data, headers)
    
    # Slide 18: Risk Management Heat Map
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SDV 리스크 관리 매트릭스"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Risk matrix grid
    grid_x, grid_y = 3, 2
    grid_width, grid_height = 10, 5
    
    # Draw grid
    for i in range(4):
        # Horizontal lines
        line_h = slide.shapes.add_connector(
            1,
            Inches(grid_x), Inches(grid_y + i * grid_height/3),
            Inches(grid_x + grid_width), Inches(grid_y + i * grid_height/3)
        )
        line_h.line.color.rgb = RGBColor(200, 200, 200)
        
        # Vertical lines
        line_v = slide.shapes.add_connector(
            1,
            Inches(grid_x + i * grid_width/3), Inches(grid_y),
            Inches(grid_x + i * grid_width/3), Inches(grid_y + grid_height)
        )
        line_v.line.color.rgb = RGBColor(200, 200, 200)
    
    # Risk items with positions (impact, probability)
    risks = [
        ("기술 격차", 2.5, 2.5, RGBColor(255, 0, 0)),  # High impact, high prob
        ("투자 부족", 2, 2, RGBColor(255, 100, 0)),
        ("인력 부족", 1.5, 2.5, RGBColor(255, 192, 0)),
        ("규제 변화", 2.5, 1, RGBColor(255, 255, 0)),
        ("시장 경쟁", 1, 2, RGBColor(100, 200, 100)),
        ("보안 위협", 3, 0.5, RGBColor(255, 0, 0))
    ]
    
    for risk_name, impact, prob, color in risks:
        x = grid_x + prob * grid_width/3
        y = grid_y + (3 - impact) * grid_height/3
        
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.4), Inches(y - 0.4),
            Inches(0.8), Inches(0.8)
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = color
        bubble.fill.transparency = 0.3
        
        # Risk label
        label_box = slide.shapes.add_textbox(
            Inches(x - 0.7), Inches(y - 0.2),
            Inches(1.4), Inches(0.4)
        )
        label_box.text_frame.text = risk_name
        label_box.text_frame.paragraphs[0].font.size = Pt(9)
        label_box.text_frame.paragraphs[0].font.bold = True
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Axis labels
    x_label = slide.shapes.add_textbox(Inches(7), Inches(7.5), Inches(2), Inches(0.5))
    x_label.text_frame.text = "발생 가능성 →"
    x_label.text_frame.paragraphs[0].font.size = Pt(12)
    x_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    y_label = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(1), Inches(2))
    y_label.text_frame.text = "영향도\n↑"
    y_label.text_frame.paragraphs[0].font.size = Pt(12)
    y_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 19: Action Plan Gantt Chart Style
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "2025년 실행 계획 (Gantt Chart)"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Tasks and timeline
    tasks = [
        ("표준화 TF 구성", 1, 2, RGBColor(0, 176, 240)),
        ("기술 격차 분석", 1, 3, RGBColor(112, 173, 71)),
        ("플랫폼 설계", 2, 4, RGBColor(255, 192, 0)),
        ("파일럿 개발", 4, 5, RGBColor(237, 125, 49)),
        ("테스트베드 구축", 6, 3, RGBColor(91, 155, 213)),
        ("파트너십 체결", 3, 9, RGBColor(112, 48, 160)),
        ("인력 채용", 1, 12, RGBColor(200, 100, 100))
    ]
    
    # Month headers
    months = ['1월', '2월', '3월', '4월', '5월', '6월', '7월', '8월', '9월', '10월', '11월', '12월']
    start_x = 3
    month_width = 0.9
    
    for i, month in enumerate(months):
        month_box = slide.shapes.add_textbox(
            Inches(start_x + i * month_width), Inches(2),
            Inches(month_width), Inches(0.4)
        )
        month_box.text_frame.text = month
        month_box.text_frame.paragraphs[0].font.size = Pt(9)
        month_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Task bars
    y_start = 2.5
    for i, (task, start_month, duration, color) in enumerate(tasks):
        # Task name
        task_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_start + i * 0.6),
            Inches(2.3), Inches(0.5)
        )
        task_box.text_frame.text = task
        task_box.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Task bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x + (start_month - 1) * month_width),
            Inches(y_start + i * 0.6),
            Inches(duration * month_width),
            Inches(0.4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
    
    # Slide 20: Investment ROI Projection
    categories = ['Y1', 'Y2', 'Y3', 'Y4', 'Y5']
    series_data = {
        '투자 (억원)': [1000, 1500, 1300, 1000, 700],
        '수익 (억원)': [200, 800, 2000, 3500, 5000],
        'ROI (%)': [20, 53, 154, 350, 714]
    }
    add_chart_slide(prs, "투자 대비 수익 전망 (5개년)", XL_CHART_TYPE.COLUMN_CLUSTERED, categories, series_data)
    
    # Slide 21: Success Metrics Summary
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "성공 지표 요약"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    
    # Create gauge-style indicators
    metrics = [
        ("기술 완성도", 85, RGBColor(112, 173, 71)),
        ("시장 준비도", 70, RGBColor(255, 192, 0)),
        ("투자 진행률", 60, RGBColor(91, 155, 213)),
        ("파트너십", 90, RGBColor(112, 173, 71))
    ]
    
    for i, (metric, value, color) in enumerate(metrics):
        x = 2 + (i % 2) * 7
        y = 2.5 + (i // 2) * 3
        
        # Background circle
        bg_circle = slide.shapes.add_shape(
            MSO_SHAPE.DONUT,
            Inches(x), Inches(y),
            Inches(2.5), Inches(2.5)
        )
        bg_circle.fill.solid()
        bg_circle.fill.fore_color.rgb = RGBColor(230, 230, 230)
        
        # Progress arc (simulated with partial donut)
        progress = slide.shapes.add_shape(
            MSO_SHAPE.DONUT,
            Inches(x), Inches(y),
            Inches(2.5), Inches(2.5)
        )
        progress.fill.solid()
        progress.fill.fore_color.rgb = color
        
        # Center text
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.5), Inches(y + 0.8),
            Inches(1.5), Inches(1)
        )
        text_frame = text_box.text_frame
        text_frame.text = f"{value}%"
        text_frame.paragraphs[0].font.size = Pt(24)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = color
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Metric label
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 2.7),
            Inches(2.5), Inches(0.5)
        )
        label_box.text_frame.text = metric
        label_box.text_frame.paragraphs[0].font.size = Pt(12)
        label_box.text_frame.paragraphs[0].font.bold = True
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 22: Conclusion with Call to Action
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background gradient effect
    bg_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(16), Inches(9)
    )
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0, 84, 159)
    bg_rect.fill.transparency = 0.1
    bg_rect.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Call to Action"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Action items with icons
    actions = [
        ("즉시 실행", ["SDV TF 구성 (2월)", "예산 확보 (3월)", "파일럿 착수 (4월)"], RGBColor(255, 0, 0)),
        ("단기 목표", ["표준화 참여", "인력 100명 확보", "테스트베드 구축"], RGBColor(255, 192, 0)),
        ("중기 목표", ["K-SDV 1.0 출시", "글로벌 파트너 10개", "시장 점유율 10%"], RGBColor(112, 173, 71)),
        ("장기 비전", ["글로벌 Top 3", "수출 10조원", "일자리 1만개"], RGBColor(0, 176, 240))
    ]
    
    y_pos = 3
    for i, (phase, items, color) in enumerate(actions):
        x_pos = 1 + i * 3.75
        
        # Phase box
        phase_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(3.5), Inches(3.5)
        )
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = color
        phase_box.fill.transparency = 0.2
        phase_box.line.color.rgb = color
        phase_box.line.width = Pt(2)
        
        # Phase title
        text_frame = phase_box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = phase
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Items
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(50, 50, 50)
    
    # Motivational quote
    quote_box = slide.shapes.add_textbox(Inches(2), Inches(7), Inches(12), Inches(1))
    quote_frame = quote_box.text_frame
    quote_frame.text = '"SDV는 선택이 아닌 필수, 미래 자동차 산업의 생존 전략"'
    quote_frame.paragraphs[0].font.size = Pt(18)
    quote_frame.paragraphs[0].font.italic = True
    quote_frame.paragraphs[0].font.color.rgb = RGBColor(0, 84, 159)
    quote_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 23: Q&A with Contact
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Background design
    bg_rect1 = slide.shapes.add_shape(
        MSO_SHAPE.WAVE,
        Inches(0), Inches(7),
        Inches(16), Inches(2)
    )
    bg_rect1.fill.solid()
    bg_rect1.fill.fore_color.rgb = RGBColor(0, 176, 240)
    bg_rect1.fill.transparency = 0.7
    bg_rect1.line.fill.background()
    
    bg_rect2 = slide.shapes.add_shape(
        MSO_SHAPE.WAVE,
        Inches(0), Inches(0),
        Inches(16), Inches(2)
    )
    bg_rect2.fill.solid()
    bg_rect2.fill.fore_color.rgb = RGBColor(0, 84, 159)
    bg_rect2.fill.transparency = 0.8
    bg_rect2.line.fill.background()
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Q&A"
    subtitle.text = "감사합니다\n\n문의사항\nsdv-korea@example.com\nTel: 02-1234-5678\n\nSDV Korea Initiative 2025"
    
    # Save the presentation
    prs.save('/home/kim/sdv-c/SDV_Advanced_Presentation_Full.pptx')
    print("Advanced presentation created successfully: SDV_Advanced_Presentation_Full.pptx")

if __name__ == "__main__":
    create_advanced_sdv_presentation()